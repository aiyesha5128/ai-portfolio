# slide_gen.py
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import os

def create_presentation(title: str, summary_text: str, chart_image_paths: list, output_path: str = "outputs/outlier_report.pptx"):
    """
    Creates a simple PPTX with a title slide, summary slide and one slide per chart image.
    """
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title

    # Summary slide
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Summary"
    body = slide.shapes.placeholders[1].text_frame
    body.text = summary_text

    # Chart slides
    for img_path in chart_image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        left = Inches(1)
        top = Inches(1.0)
        slide.shapes.add_picture(img_path, left, top, width=Inches(8))

    prs.save(output_path)
    return output_path
